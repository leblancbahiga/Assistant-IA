"""
NURU V10 — Spotlight Search V2 : recherche macOS + lecture + scoring.

Améliorations V2 :
  - mdfind avec tous les termes (pas limité à 3)
  - Recherche plein texte (kMDItemTextContent) avec préfixe astring
  - Scoring de pertinence basé sur : nb de termes trouvés dans le chemin,
    dans le nom de fichier, ratio dans le contenu
  - Historique des recherches
  - Exclusion robuste du projet NURU
"""
import subprocess
import logging
import os
import re
from dataclasses import dataclass
from typing import List, Optional
from pathlib import Path

logger = logging.getLogger(__name__)

# Extensions lisibles (ordre de pertinence décroissant)
READABLE_EXTS = {".txt", ".md", ".py", ".csv", ".json", ".html", ".rtf",
                 ".docx", ".pdf", ".pptx", ".xlsx", ".yaml", ".yml", ".toml",
                 ".ini", ".cfg", ".log", ".xml", ".js", ".ts", ".sql"}

# Mots vides (français + anglais)
STOP_WORDS = {
    "le", "la", "les", "de", "du", "des", "un", "une", "et", "ou",
    "est", "sont", "a", "ai", "as", "avons", "avez", "ont",
    "parle", "moi", "dit", "raconte", "montre", "donne",
    "quel", "quelle", "quels", "quelles", "que", "qui",
    "comment", "pourquoi", "quand", "où", "combien",
    "the", "a", "an", "and", "or", "is", "are", "was", "were",
}


@dataclass
class SpotlightResult:
    path: str
    filename: str
    content: str = ""
    relevance: float = 0.0
    match_count: int = 0       # Nb de termes distincts trouvés
    content_match: bool = False


def _is_project_path(path: str) -> bool:
    """Vrai si le chemin fait partie du projet NURU."""
    if "Assistant IA" not in path:
        return False
    blocked_subdirs = [
        "/src/", "/tests/", "/scripts/", "/.venv/", "/docs/",
        "/.git/", "/node_modules/", "/__pycache__/",
    ]
    return any(part in path for part in blocked_subdirs)


class SpotlightSearch:
    """Recherche via Spotlight (mdfind) avec lecture + scoring V2."""

    def __init__(self):
        self._is_macos = os.uname().sysname == "Darwin"
        self._search_history: list = []
        if not self._is_macos:
            logger.warning("Spotlight non disponible (pas macOS)")

    # ── Extraction ──────────────────────────────────────────────────────

    def _extract_key_terms(self, query: str) -> List[str]:
        """Extrait les termes significatifs d'une requête."""
        words = re.findall(r"\w{2,}", query.lower())
        return [w for w in words if w not in STOP_WORDS]

    # ── Recherche ───────────────────────────────────────────────────────

    def search(self, query: str, scope: Optional[str] = None,
               max_results: int = 10, read_content: bool = True,
               content_search: bool = True) -> List[SpotlightResult]:
        """Cherche via Spotlight avec plein texte en option.

        Args:
            query: Requête textuelle
            scope: Dossier racine (None = partout)
            max_results: Nombre max de résultats
            read_content: Lire le contenu des fichiers trouvés
            content_search: Activer la recherche plein texte (astring)
        """
        if not self._is_macos:
            return []

        key_terms = self._extract_key_terms(query)
        if not key_terms:
            return []

        # 1) Recherche dans les noms de fichiers
        results = self._mdfind(key_terms, scope, max_results,
                               content_search=False)
        # 2) Recherche dans le contenu (complémentaire)
        if content_search:
            content_results = self._mdfind(key_terms, scope,
                                           max_results,
                                           content_search=True)
            # Fusion des résultats (les content_results n'ont pas
            # forcement le contenu lu)
            existing_paths = {r.path for r in results}
            for cr in content_results:
                if cr.path not in existing_paths:
                    results.append(cr)
                    existing_paths.add(cr.path)

        results.sort(key=lambda r: r.relevance, reverse=True)

        # Lecture du contenu si demandé
        if read_content:
            for r in results[:max_results]:
                if not r.content:
                    r.content = self._read_file(Path(r.path), max_chars=2000)

        final = results[:max_results]
        self._search_history.append({"query": query,
                                     "results": len(final)})
        return final

    # ── mdfind ──────────────────────────────────────────────────────────

    def _mdfind(self, key_terms: List[str], scope: Optional[str],
                max_results: int,
                content_search: bool) -> List[SpotlightResult]:
        """Exécute mdfind et parse les résultats."""
        cmd = ["mdfind"]
        if scope:
            scope_path = Path(scope).expanduser()
            if scope_path.exists():
                cmd.extend(["-onlyin", str(scope_path)])

        if content_search:
            # Recherche plein texte via kMDItemTextContent
            # Préfixe "astring" pour recherche-insensible dans tout le contenu
            query_parts = [f'kMDItemTextContent == "*{t}*"c'
                           for t in key_terms]
            cmd.append(query_parts[0] if len(query_parts) == 1
                       else f"({' && '.join(query_parts)})")
        else:
            # Recherche classique mots-clés sur le nom + chemin
            search_query = " ".join(key_terms[:5])
            cmd.append(search_query)

        try:
            result = subprocess.run(cmd, capture_output=True, text=True,
                                    timeout=15)
            if result.returncode != 0:
                return []

            paths = [p.strip()
                     for p in result.stdout.strip().split("\n")
                     if p.strip()]
            return self._parse_results(paths, key_terms, max_results,
                                       content_search)

        except subprocess.TimeoutExpired:
            logger.warning("Spotlight timeout: %s", " ".join(key_terms))
            return []
        except Exception as e:
            logger.error("Erreur Spotlight: %s", e)
            return []

    # ── Parsing / Scoring ───────────────────────────────────────────────

    def _parse_results(self, paths: List[str], key_terms: List[str],
                       max_results: int,
                       content_match: bool) -> List[SpotlightResult]:
        """Parse les chemins, calcule la pertinence, filtre."""
        results = []
        for path in paths[:max_results * 5]:
            p = Path(path)
            if not p.exists() or not p.is_file():
                continue
            if _is_project_path(path):
                continue
            if p.suffix.lower() not in READABLE_EXTS:
                continue

            # Scoring de pertinence
            name_lower = p.name.lower()
            path_lower = path.lower()
            match_count = 0
            relevance = 0.0

            for term in key_terms:
                if term in name_lower:
                    match_count += 1
                    relevance += 2.0  # Bonus fort pour nom de fichier
                elif term in path_lower:
                    match_count += 1
                    relevance += 1.0  # Bonus moyen pour chemin

            # Bonus si la recherche a matché le contenu
            if content_match:
                relevance += 1.5

            # Normalisation
            relevance = relevance / max(len(key_terms), 1)

            # Contenu (lecture différée si read_content=True)
            content = ""
            if match_count > 0:
                content = self._read_file(p, max_chars=500)

            results.append(SpotlightResult(
                path=str(p),
                filename=p.name,
                content=content,
                relevance=min(relevance, 1.0),
                match_count=match_count,
                content_match=content_match,
            ))

        results.sort(key=lambda r: (r.match_count, r.relevance),
                     reverse=True)
        return results[:max_results]

    # ── Lecture fichiers ────────────────────────────────────────────────

    def _read_file(self, path: Path, max_chars: int = 2000) -> str:
        """Lit le contenu d'un fichier selon son type."""
        suffix = path.suffix.lower()
        try:
            if suffix in {".txt", ".md", ".py", ".csv", ".json", ".html",
                          ".yaml", ".yml", ".toml", ".ini", ".cfg", ".log",
                          ".xml", ".js", ".ts", ".sql"}:
                text = path.read_text(encoding="utf-8", errors="ignore")
                return text[:max_chars]
            elif suffix == ".pdf":
                import pymupdf
                doc = pymupdf.open(str(path))
                text = "".join(page.get_text() for page in doc)
                doc.close()
                return text[:max_chars]
            elif suffix == ".docx":
                from docx import Document
                doc = Document(str(path))
                text = "\n".join(p.text for p in doc.paragraphs)
                return text[:max_chars]
            elif suffix == ".pptx":
                from pptx import Presentation
                prs = Presentation(str(path))
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                    if len(text) > max_chars:
                        break
                return text[:max_chars]
            elif suffix == ".xlsx":
                from openpyxl import load_workbook
                wb = load_workbook(str(path), data_only=True)
                text = ""
                for sheet in wb.sheetnames:
                    ws = wb[sheet]
                    for row in ws.iter_rows(values_only=True):
                        text += " ".join(str(c) for c in row
                                         if c is not None) + "\n"
                    if len(text) > max_chars:
                        break
                return text[:max_chars]
            return ""
        except Exception as e:
            logger.debug("Erreur lecture %s: %s", path.name, e)
            return ""

    def get_history(self) -> list:
        return list(self._search_history)

    def clear_history(self) -> None:
        self._search_history.clear()
