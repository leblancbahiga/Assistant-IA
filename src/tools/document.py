from __future__ import annotations
from dataclasses import dataclass, field
from enum import Enum
from pathlib import Path
from typing import Any


class DocFormat(str, Enum):
    WORD = "word"
    PDF = "pdf"
    PPTX = "pptx"
    XLSX = "xlsx"
    MARKDOWN = "markdown"


@dataclass
class DocSection:
    title: str
    content: str
    level: int = 1  # h1=1, h2=2...


@dataclass
class DocumentSpec:
    title: str
    format: DocFormat
    sections: list[DocSection]
    metadata: dict[str, str] = field(default_factory=dict)


class DocumentGenerator:
    """Génère de vrais fichiers Word, PDF, PPTX, XLSX et Markdown."""

    SUPPORTED_FORMATS = set(DocFormat)

    # ── Validation ────────────────────────────────────────────

    def validate_spec(self, spec: DocumentSpec) -> list[str]:
        """Valide un DocumentSpec. Retourne les erreurs (vide = OK)."""
        errors: list[str] = []
        if not spec.title:
            errors.append("Titre requis")
        if not spec.sections:
            errors.append("Au moins une section requise")
        for i, s in enumerate(spec.sections):
            if not s.title:
                errors.append(f"Section {i + 1}: titre requis")
            if not s.content:
                errors.append(f"Section {i + 1}: contenu requis")
        if spec.format not in self.SUPPORTED_FORMATS:
            errors.append(f"Format non supporté: {spec.format}")
        return errors

    # ── Routeur principal ─────────────────────────────────────

    def generate(self, spec: DocumentSpec, output_path: str | Path) -> Path:
        """Génère le fichier dans le format demandé. Retourne le chemin du fichier."""
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        dispatch = {
            DocFormat.WORD: self._generate_word,
            DocFormat.PDF: self._generate_pdf,
            DocFormat.PPTX: self._generate_pptx,
            DocFormat.XLSX: self._generate_xlsx,
            DocFormat.MARKDOWN: self._generate_markdown,
        }
        handler = dispatch.get(spec.format)
        if handler is None:
            raise ValueError(f"Format non supporté: {spec.format}")
        return handler(spec, output_path)

    # ── Word (.docx) ──────────────────────────────────────────

    def _generate_word(self, spec: DocumentSpec, path: Path) -> Path:
        from docx import Document
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.shared import Pt

        doc = Document()
        # Titre
        title_para = doc.add_heading(spec.title, level=0)
        title_para.alignment = WD_ALIGN_PARAGRAPH.CENTER
        # Métadonnées
        if spec.metadata:
            for k, v in spec.metadata.items():
                p = doc.add_paragraph()
                run = p.add_run(f"{k}: {v}")
                run.font.size = Pt(10)
                run.font.italic = True
            doc.add_paragraph()  # espace
        # Sections
        for section in spec.sections:
            level = min(section.level, 4)
            doc.add_heading(section.title, level=level)
            doc.add_paragraph(section.content)
        doc.save(str(path))
        return path

    # ── PDF ───────────────────────────────────────────────────

    def _generate_pdf(self, spec: DocumentSpec, path: Path) -> Path:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.units import cm
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

        doc = SimpleDocTemplate(
            str(path),
            pagesize=A4,
            topMargin=2 * cm,
            bottomMargin=2 * cm,
            leftMargin=2 * cm,
            rightMargin=2 * cm,
        )
        styles = getSampleStyleSheet()
        story: list = []

        # Titre
        story.append(Paragraph(spec.title, styles["Title"]))
        story.append(Spacer(1, 12))
        # Métadonnées
        for k, v in spec.metadata.items():
            story.append(Paragraph(f"<i>{k}: {v}</i>", styles["Normal"]))
        if spec.metadata:
            story.append(Spacer(1, 12))
        # Sections
        for section in spec.sections:
            heading_key = f"Heading{min(section.level, 4)}"
            style = styles.get(heading_key, styles["Heading2"])
            story.append(Paragraph(section.title, style))
            story.append(Paragraph(section.content, styles["Normal"]))
            story.append(Spacer(1, 8))

        doc.build(story)
        return path

    # ── PowerPoint (.pptx) ────────────────────────────────────

    def _generate_pptx(self, spec: DocumentSpec, path: Path) -> Path:
        from pptx import Presentation

        prs = Presentation()
        # Slide titre
        slide_layout = prs.slide_layouts[0]  # Title Slide
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = spec.title
        if spec.metadata:
            slide.placeholders[1].text = "\n".join(
                f"{k}: {v}" for k, v in spec.metadata.items()
            )
        # Slides sections
        for section in spec.sections:
            slide_layout = prs.slide_layouts[1]  # Title and Content
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = section.title
            slide.placeholders[1].text = section.content
        prs.save(str(path))
        return path

    # ── Excel (.xlsx) ─────────────────────────────────────────

    def _generate_xlsx(self, spec: DocumentSpec, path: Path) -> Path:
        from openpyxl import Workbook
        from openpyxl.styles import Alignment, Font

        wb = Workbook()
        ws = wb.active
        ws.title = spec.title[:31]  # max 31 chars for sheet name
        # Titre
        ws.merge_cells("A1:D1")
        title_cell = ws["A1"]
        title_cell.value = spec.title
        title_cell.font = Font(bold=True, size=14)
        title_cell.alignment = Alignment(horizontal="center")
        # Métadonnées
        row = 3
        for k, v in spec.metadata.items():
            ws.cell(row=row, column=1, value=k).font = Font(bold=True)
            ws.cell(row=row, column=2, value=v)
            row += 1
        row += 1  # espace
        # En-têtes
        ws.cell(row=row, column=1, value="Section").font = Font(bold=True)
        ws.cell(row=row, column=2, value="Contenu").font = Font(bold=True)
        row += 1
        # Lignes
        for section in spec.sections:
            ws.cell(row=row, column=1, value=section.title)
            ws.cell(row=row, column=2, value=section.content)
            row += 1
        # Ajuster largeurs
        ws.column_dimensions["A"].width = 30
        ws.column_dimensions["B"].width = 60
        wb.save(str(path))
        return path

    # ── Markdown ──────────────────────────────────────────────

    def _generate_markdown(self, spec: DocumentSpec, path: Path) -> Path:
        lines = [f"# {spec.title}", ""]
        if spec.metadata:
            for k, v in spec.metadata.items():
                lines.append(f"**{k}** : {v}")
            lines.append("")
        for section in spec.sections:
            prefix = "#" * min(section.level + 1, 6)
            lines.extend([f"{prefix} {section.title}", "", section.content, ""])
        path.write_text("\n".join(lines), encoding="utf-8")
        return path

    # ── Utilitaires ───────────────────────────────────────────

    def generate_markdown_string(self, spec: DocumentSpec) -> str:
        """Génère le contenu Markdown en string (sans écrire de fichier)."""
        lines = [f"# {spec.title}", ""]
        if spec.metadata:
            for k, v in spec.metadata.items():
                lines.append(f"**{k}** : {v}")
            lines.append("")
        for section in spec.sections:
            prefix = "#" * min(section.level + 1, 6)
            lines.extend([f"{prefix} {section.title}", "", section.content, ""])
        return "\n".join(lines)

    def generate_json(self, spec: DocumentSpec) -> dict:
        """Génère une représentation JSON du document."""
        return {
            "title": spec.title,
            "format": spec.format.value,
            "metadata": spec.metadata,
            "sections": [
                {"title": s.title, "content": s.content, "level": s.level}
                for s in spec.sections
            ],
        }

    def generate_from_template(self, template: str, data: dict[str, str]) -> str:
        """Remplace les placeholders {{key}} dans un template."""
        result = template
        for key, value in data.items():
            result = result.replace("{{" + key + "}}", str(value))
        return result
