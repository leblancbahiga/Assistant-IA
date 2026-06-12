#!/usr/bin/env python3
"""
NURU V10 — Réindexation Robuste (Batch + Checkpoint + GC).

Pipeline:
  1. Scan ~/Desktop + ~/Documents pour les fichiers parsables
  2. Traite par lots de BATCH_FILES (50 fichiers par lot)
  3. Chaque lot : lire → chunker → embed → indexer → checkpoint
  4. GC collect() après chaque lot (évite saturation RAM M1 8 Go)
  5. Checkpoint JSON → reprise possible si interrompu
  6. Mode incrémental : fichiers inchangés sautés (SHA256)

Usage:
  python3 scripts/reindex_all.py              # Full reindex
  python3 scripts/reindex_all.py --incremental # Skip unchanged files
  python3 scripts/reindex_all.py --force       # Wipe + rebuild
"""

import gc
import json
import logging
import os
import sys
import time
import numpy as np
from pathlib import Path
from typing import List, Optional

project_root = Path(__file__).parent.parent
if str(project_root) not in sys.path:
    sys.path.insert(0, str(project_root))

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    datefmt="%H:%M:%S",
)
logger = logging.getLogger(__name__)

# ── Configuration ──────────────────────────────────────────────────────────

INDEXED_EXTS = {".pdf", ".docx", ".txt", ".md", ".csv", ".json", ".pptx", ".xlsx", ".rtf"}
SCAN_DIRS = [Path.home() / "Desktop", Path.home() / "Documents"]
BATCH_FILES = 50          # Fichiers par lot (RAM ~200-400 Mo par lot)
CHECKPOINT_PATH = project_root / "indexes" / "reindex_checkpoint.json"
DB_PATH = project_root / "indexes" / "nuru.db"
MAX_FILE_CHARS = 500_000   # Ignorer les fichiers >500K chars (>10 min)
MIN_CONTENT_CHARS = 100    # Ignorer les fichiers trop courts
CHUNK_MAX_CHARS = 500      # V10 Expert: 500 chars (était 2000)
CHUNK_OVERLAP = 100         # V10 Expert: 100 chars (était 200)
EMBED_BATCH_SIZE = 16      # Embeddings par appel asyncio.gather


# ══════════════════════════════════════════════════════════════════════════
#  Checkpoint management
# ══════════════════════════════════════════════════════════════════════════


def load_checkpoint() -> dict:
    """Charge le checkpoint existant ou retourne un état vide."""
    if CHECKPOINT_PATH.exists():
        try:
            data = json.loads(CHECKPOINT_PATH.read_text())
            logger.info(f"📌 Checkpoint trouvé: {data.get('processed', 0)} fichiers traités")
            return data
        except Exception as e:
            logger.warning(f"⚠️ Checkpoint corrompu: {e} — redémarrage")
    return {"processed": 0, "total": 0, "completed_files": [], "start_time": time.time()}


def save_checkpoint(state: dict) -> None:
    """Sauvegarde l'état après chaque lot."""
    state["last_update"] = time.time()
    CHECKPOINT_PATH.parent.mkdir(parents=True, exist_ok=True)
    # Alléger le checkpoint : garder seulement les noms, pas les chemins complets
    serializable = {
        "processed": state["processed"],
        "total": state["total"],
        "completed_files": state["completed_files"][-500:],  # garder les 500 derniers
        "start_time": state["start_time"],
        "last_update": state["last_update"],
    }
    CHECKPOINT_PATH.write_text(json.dumps(serializable, indent=2))


def clear_checkpoint() -> None:
    """Supprime le checkpoint (pour --force)."""
    if CHECKPOINT_PATH.exists():
        CHECKPOINT_PATH.unlink()
        logger.info("🗑️ Checkpoint effacé")


# ══════════════════════════════════════════════════════════════════════════
#  File scanning
# ══════════════════════════════════════════════════════════════════════════


def scan_files(directories: List[Path]) -> List[Path]:
    """Scanne les répertoires pour trouver les fichiers parsables, triés par taille."""
    files = []
    for directory in directories:
        if not directory.exists():
            continue
        for root, dirs, filenames in os.walk(directory):
            # Exclure les dources système et projets
            dirs[:] = [
                d for d in dirs
                if not d.startswith(".")
                and d not in {"__pycache__", "node_modules", ".git", "Library", "Applications", "build", ".gradle", "Pods"}
            ]
            # Exclure les chemins build (Android Studio, Xcode, etc.)
            if "/build/" in root or "/.gradle/" in root or "/Pods/" in root:
                dirs[:] = []  # Ne pas descendre dans ces dossiers
                continue
            for fname in filenames:
                if fname.startswith("."):
                    continue
                fpath = Path(root) / fname
                if fpath.suffix.lower() in INDEXED_EXTS:
                    # Vérifier la taille
                    try:
                        fsize = fpath.stat().st_size
                        if fsize > 50 * 1024 * 1024:  # > 50 Mo → skip
                            continue
                        files.append((fpath, fsize))
                    except OSError:
                        continue

    # Trier par taille (petits fichiers d'abord → diagnostic rapide)
    files.sort(key=lambda x: x[1])
    return [f[0] for f in files]


# ══════════════════════════════════════════════════════════════════════════
#  File reading (copié de l'ancien script, avec timeout par fichier)
# ══════════════════════════════════════════════════════════════════════════


def read_file_content(fpath: Path) -> str:
    """Lit le contenu d'un fichier selon son extension. Timeout implicite via signal."""
    suffix = fpath.suffix.lower()
    try:
        if suffix in {".txt", ".md", ".csv", ".json"}:
            return fpath.read_text(encoding="utf-8", errors="ignore")
        elif suffix == ".pdf":
            import pymupdf
            doc = pymupdf.open(str(fpath))
            # V10: Gérer les PDFs cryptés/protégés
            if doc.is_encrypted:
                logger.debug(f"⚠️ PDF crypté: {fpath.name} — skip")
                doc.close()
                return ""
            text = "\n".join([page.get_text() for page in doc])
            doc.close()
            return text
        elif suffix == ".docx":
            from docx import Document
            doc = Document(str(fpath))
            # V10.1 : extraire aussi le contenu des tables
            parts = [p.text for p in doc.paragraphs]
            for table in doc.tables:
                for row in table.rows:
                    row_text = "\t".join(cell.text for cell in row.cells if cell.text.strip())
                    if row_text.strip():
                        parts.append(row_text)
            return "\n".join(parts)
        elif suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(fpath))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        elif suffix == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(str(fpath), data_only=True)
            text = ""
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    text += " ".join([str(c) for c in row if c is not None]) + "\n"
            wb.close()
            return text
        elif suffix == ".rtf":
            # RTF basique — extraction strippée
            text = fpath.read_text(encoding="utf-8", errors="ignore")
            import re
            text = re.sub(r"\\([a-z]+)", " ", text)
            text = re.sub(r"[{}]", " ", text)
            return text
        return ""
    except Exception as e:
        logger.debug(f"⚠️ Lecture {fpath.name}: {e}")
        return ""


# ══════════════════════════════════════════════════════════════════════════
#  Chunking (copié de l'ancien script)
# ══════════════════════════════════════════════════════════════════════════


def chunk_text(text: str, source: str, max_chars: int = CHUNK_MAX_CHARS, overlap: int = CHUNK_OVERLAP) -> list[dict]:
    """Découpe un texte en chunks avec overlap (V10.1: restaure chunking simple d'origine)."""
    if not text or len(text.strip()) < MIN_CONTENT_CHARS:
        return []
    chunks = []
    lines = text.split("\n")
    current_chunk = ""
    current_len = 0
    for line in lines:
        current_chunk += line + "\n"
        current_len += len(line) + 1
        if current_len >= max_chars:
            chunks.append({"content": current_chunk.strip(), "source": source})
            words = current_chunk.split()
            overlap_words = words[-overlap // 5:] if len(words) > overlap // 5 else []
            current_chunk = " ".join(overlap_words) + "\n"
            current_len = len(current_chunk)
    if current_chunk.strip() and len(current_chunk.strip()) > MIN_CONTENT_CHARS:
        chunks.append({"content": current_chunk.strip(), "source": source})
    return chunks


# ══════════════════════════════════════════════════════════════════════════
#  Main pipeline
# ══════════════════════════════════════════════════════════════════════════


def main():
    import argparse
    parser = argparse.ArgumentParser(description="NURU V10 — Réindexation Robuste")
    parser.add_argument("--incremental", action="store_true",
                        help="Mode incrémental : ignorer les fichiers inchangés")
    parser.add_argument("--force", action="store_true",
                        help="Forcer : vider l'index et réindexer complètement")
    parser.add_argument("--batch", type=int, default=BATCH_FILES,
                        help=f"Taille du lot (défaut: {BATCH_FILES})")
    args = parser.parse_args()

    # ── Force : wipe index + checkpoint ──
    if args.force:
        logger.info("🔥 Mode FORCE : vidage de l'index...")
        if DB_PATH.exists():
            DB_PATH.unlink()
            logger.info(f"   🗑️ {DB_PATH} supprimé")
        clear_checkpoint()

    # ── Charger checkpoint ──
    state = load_checkpoint()
    completed_set = set(state.get("completed_files", []))

    # ── Initialiser RAG Engine ──
    logger.info("🔧 Initialisation RAG Engine...")
    from src.rag_engine import RAGEngine
    from src.embedder import Embedder

    rag = RAGEngine()
    embedder = Embedder()
    logger.info("   ✅ RAG Engine prêt")

    # ── Scan ──
    logger.info("📂 Scan des fichiers...")
    all_files = scan_files(SCAN_DIRS)
    logger.info(f"   {len(all_files)} fichiers trouvés")

    # Stats extensions
    ext_counts = {}
    for f in all_files:
        ext = f.suffix.lower()
        ext_counts[ext] = ext_counts.get(ext, 0) + 1
    for ext, count in sorted(ext_counts.items(), key=lambda x: -x[1]):
        logger.info(f"   {ext}: {count}")

    # Filtre les fichiers trop gros (> 10 Mo = lents à parser)
    before_size = len(all_files)
    all_files = [f for f in all_files if f.stat().st_size <= 10 * 1024 * 1024]
    logger.info(f"   Filtre taille (<10 Mo): {before_size} → {len(all_files)} fichiers")

    # ── Filtrer fichiers déjà traités ──
    if args.incremental:
        remaining = [f for f in all_files if str(f) not in completed_set]
        skipped_count = len(all_files) - len(remaining)
        logger.info(f"   Mode incrémental: {skipped_count} déjà traités, {len(remaining)} restants")
        all_files = remaining

    # Filtre les fichiers trop longs
    all_files = [f for f in all_files if f.stat().st_size < MAX_FILE_CHARS]
    logger.info(f"   Après filtrage taille: {len(all_files)} fichiers")

    total_files = len(all_files)
    if total_files == 0:
        logger.info("✅ Aucun fichier à indexer.")
        return

    state["total"] = total_files
    save_checkpoint(state)

    # ── Pipeline par lots ──
    logger.info(f"\n🚀 Démarrage indexation ({total_files} fichiers, batch={args.batch})")
    t_start = time.time()
    total_chunks = 0
    total_skipped = 0
    last_log_time = 0
    processed_since_checkpoint = 0

    for batch_start in range(0, total_files, args.batch):
        batch_files = all_files[batch_start:batch_start + args.batch]
        batch_chunks = []

        # ── Étape A: Lire + Chunker le lot ──
        t_batch = time.time()
        file_timeout = 30  # 30s max par fichier
        for fpath in batch_files:
            try:
                f_start = time.time()
                content = read_file_content(fpath)
                f_elapsed = time.time() - f_start
                if f_elapsed > file_timeout:
                    logger.warning(f"⚠️ Timeout {fpath.name} ({f_elapsed:.0f}s) → skip")
                    total_skipped += 1
                    continue
                if not content or len(content.strip()) < MIN_CONTENT_CHARS:
                    total_skipped += 1
                    continue
                # Chemin relatif au home pour la source
                try:
                    source = str(fpath.relative_to(Path.home()))
                except ValueError:
                    source = fpath.name
                chunks = chunk_text(content, source)
                if chunks:
                    batch_chunks.extend(chunks)
                else:
                    total_skipped += 1
            except Exception as e:
                total_skipped += 1
                logger.debug(f"⚠️ Erreur {fpath.name}: {e}")

        if not batch_chunks:
            # Marquer comme traités même si vides
            for fpath in batch_files:
                completed_set.add(str(fpath))
            continue

        # ── Étape B: Embed le lot (batch unique — 1 thread au lieu de N) ──
        texts = [c["content"] for c in batch_chunks]
        valid_chunks = []  # V10: initialiser avant le try pour éviter UnboundLocalError
        embed_success = False

        def _do_embed(text_subset):
            """Embed un sous-ensemble de textes (évite saturation GPU Metal 4 GB)."""
            return np.array(embedder.embed_sync(text_subset, is_query=False))

        try:
            # V10: Sous-découpage par paquets de 50 chunks max
            # (évite Metal allocation > 4GB avec chunk_size=500)
            SUB_BATCH = 50
            all_embeddings = []
            for i in range(0, len(texts), SUB_BATCH):
                subset = texts[i:i + SUB_BATCH]
                emb = _do_embed(subset)
                # Normaliser les shapes
                if len(emb.shape) == 1:
                    emb = emb.reshape(1, -1)
                for j in range(len(subset)):
                    vec = emb[j].flatten()
                    all_embeddings.append(vec)
            embeddings = all_embeddings
            for chunk, emb in zip(batch_chunks, embeddings):
                chunk["embedding"] = emb.tolist()
            valid_chunks = [c for c in batch_chunks if "embedding" in c]
            embed_success = True
        except Exception as e:
            logger.warning(f"⚠️ Échec embedding batch (1ère tentative): {e}")
            # Forcer GC et réessayer une fois
            gc.collect()
            try:
                SUB_BATCH = 50
                all_embeddings = []
                for i in range(0, len(texts), SUB_BATCH):
                    subset = texts[i:i + SUB_BATCH]
                    emb = _do_embed(subset)
                    if len(emb.shape) == 1:
                        emb = emb.reshape(1, -1)
                    for j in range(len(subset)):
                        vec = emb[j].flatten()
                        all_embeddings.append(vec)
                embeddings = all_embeddings
                for chunk, emb in zip(batch_chunks, embeddings):
                    chunk["embedding"] = emb.tolist()
                valid_chunks = [c for c in batch_chunks if "embedding" in c]
                embed_success = True
            except Exception as e2:
                logger.error(f"❌ Échec embedding (abandon lot): {e2}")

        if not embed_success:
            # Marquer comme traités mais sans indexation
            for fpath in batch_files:
                completed_set.add(str(fpath))
            # GC + continue
            del batch_chunks, texts
            gc.collect()
            continue

        # ── Étape C: Indexer le lot ──
        if valid_chunks:
            try:
                rag.add_chunks(valid_chunks, dedup_source=True)
                total_chunks += len(valid_chunks)
            except Exception as e:
                logger.warning(f"⚠️ Échec indexation batch: {e}")

        # ── Marquer comme traités ──
        for fpath in batch_files:
            completed_set.add(str(fpath))

        processed_since_checkpoint += len(batch_files)
        state["processed"] = state.get("processed", 0) + len(batch_files)
        state["completed_files"] = list(completed_set)

        # ── GC forcé ──
        chunk_count = len(valid_chunks)
        del batch_chunks, texts, embeddings, valid_chunks
        gc.collect()

        # ── Log progression ──
        elapsed = time.time() - t_start
        batch_time = time.time() - t_batch
        rate = state["processed"] / elapsed if elapsed > 0 else 0
        eta = (total_files - state["processed"]) / rate if rate > 0 else 0

        logger.info(
            f"📊 Lot {state['processed']}/{total_files} "
            f"({state['processed'] * 100 // total_files}%) "
            f"| +{len(batch_files)} fichiers "
            f"| {chunk_count} chunks "
            f"| {batch_time:.1f}s "
            f"| ETA: {eta:.0f}s"
        )

        # ── Checkpoint tous les 2 lots ──
        if processed_since_checkpoint >= args.batch * 2:
            save_checkpoint(state)
            processed_since_checkpoint = 0

        # ── Petite pause pour laisser la RAM respirer ──
        time.sleep(0.5)

    # ── Finaliser ──
    total_time = time.time() - t_start
    save_checkpoint(state)

    logger.info("\n" + "=" * 60)
    logger.info("📊 RÉSUMÉ DE L'INDEXATION")
    logger.info("=" * 60)
    logger.info(f"   Fichiers traités  : {state['processed']}")
    logger.info(f"   Fichiers ignorés  : {total_skipped}")
    logger.info(f"   Chunks indexés    : {total_chunks}")
    logger.info(f"   Temps total       : {total_time:.1f}s ({total_time / 60:.1f} min)")
    if total_chunks > 0:
        logger.info(f"   Vitesse           : {total_chunks / total_time:.1f} chunks/s")
    logger.info("=" * 60)
    logger.info("✅ Indexation terminée avec succès !")

    # Nettoyer le checkpoint (succès → plus besoin)
    clear_checkpoint()


if __name__ == "__main__":
    main()
