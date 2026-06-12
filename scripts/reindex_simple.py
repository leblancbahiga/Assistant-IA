#!/usr/bin/env python3
"""
NURU V10 — Réindexation simplifiée.
Indexe les fichiers de ~/Desktop et ~/Documents.
"""

import os
import sys
import time
import asyncio
import logging
import numpy as np
from pathlib import Path
from typing import List, Dict

project_root = Path(__file__).parent.parent
if str(project_root) not in sys.path:
    sys.path.insert(0, str(project_root))

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
logger = logging.getLogger(__name__)

INDEXED_EXTS = {".pdf", ".docx", ".txt", ".md", ".csv", ".json", ".py", ".pptx", ".html", ".xlsx"}
SCAN_DIRS = [Path.home() / "Desktop", Path.home() / "Documents"]


def scan_files(directories: List[Path]) -> List[Path]:
    files = []
    for directory in directories:
        if not directory.exists():
            continue
        for root, dirs, filenames in os.walk(directory):
            dirs[:] = [d for d in dirs if not d.startswith('.') and d not in {'__pycache__', 'node_modules', '.git'}]
            for fname in filenames:
                if fname.startswith('.'):
                    continue
                fpath = Path(root) / fname
                if fpath.suffix.lower() in INDEXED_EXTS:
                    files.append(fpath)
    return files


def read_file_content(fpath: Path) -> str:
    suffix = fpath.suffix.lower()
    try:
        if suffix in {".txt", ".md", ".py", ".csv", ".json", ".html"}:
            return fpath.read_text(encoding="utf-8", errors="ignore")
        elif suffix == ".pdf":
            import pymupdf
            doc = pymupdf.open(str(fpath))
            text = "\n".join([page.get_text() for page in doc])
            doc.close()
            return text
        elif suffix == ".docx":
            from docx import Document
            doc = Document(str(fpath))
            return "\n".join([p.text for p in doc.paragraphs])
        elif suffix == ".pptx":
            from pptx import Presentation
            prs = Presentation(str(fpath))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
            return text
        elif suffix == ".xlsx":
            from openpyxl import load_workbook
            wb = load_workbook(str(fpath), data_only=True)
            text = ""
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                for row in ws.iter_rows(values_only=True):
                    text += " ".join([str(c) for c in row if c is not None]) + "\n"
            return text
        return ""
    except Exception as e:
        return ""


def chunk_text(text: str, source: str, max_chars: int = 2000, overlap: int = 200) -> List[Dict]:
    if not text or len(text.strip()) < 100:
        return []
    chunks = []
    lines = text.split("\n")
    current_chunk = ""
    current_len = 0
    for line in lines:
        current_chunk += line + "\n"
        current_len += len(line) + 1
        if current_len >= max_chars:
            chunks.append({"content": current_chunk.strip(), "source": source})
            words = current_chunk.split()
            overlap_words = words[-overlap // 5:] if len(words) > overlap // 5 else []
            current_chunk = " ".join(overlap_words) + "\n"
            current_len = len(current_chunk)
    if current_chunk.strip() and len(current_chunk.strip()) > 100:
        chunks.append({"content": current_chunk.strip(), "source": source})
    return chunks


def main():
    logger.info("=" * 60)
    logger.info("NURU V10 — Réindexation simplifiée")
    logger.info("=" * 60)

    # 1. Scanner
    logger.info("📂 Scan...")
    all_files = scan_files(SCAN_DIRS)
    logger.info(f"   {len(all_files)} fichiers trouvés")

    ext_counts = {}
    for f in all_files:
        ext = f.suffix.lower()
        ext_counts[ext] = ext_counts.get(ext, 0) + 1
    for ext, count in sorted(ext_counts.items(), key=lambda x: -x[1]):
        logger.info(f"   {ext}: {count}")

    # 2. Initialiser
    logger.info("\n🔧 Initialisation...")
    from src.rag_engine import RAGEngine
    from src.embedder import Embedder
    rag = RAGEngine()
    embedder = Embedder()
    logger.info("   ✅ OK")

    # 3. Lire et indexer
    logger.info("\n📝 Indexation...")
    total_chunks = 0
    total_files = 0
    skipped = 0
    errors = 0
    start = time.time()

    for i, fpath in enumerate(all_files):
        if (i + 1) % 50 == 0:
            logger.info(f"   [{i+1}/{len(all_files)}] {fpath.name}...")

        try:
            content = read_file_content(fpath)
            if not content or len(content.strip()) < 100:
                skipped += 1
                continue

            source = str(fpath)
            chunks = chunk_text(content, source)
            if not chunks:
                skipped += 1
                continue

            # Embed
            for chunk in chunks:
                try:
                    embedding = asyncio.run(embedder.embed(chunk["content"]))
                    if hasattr(embedding, 'shape') and len(embedding.shape) > 1:
                        embedding = embedding.flatten()
                    chunk["embedding"] = embedding.tolist() if hasattr(embedding, 'tolist') else list(embedding)
                except Exception as e:
                    continue

            valid_chunks = [c for c in chunks if "embedding" in c]
            if not valid_chunks:
                skipped += 1
                continue

            rag.add_chunks(valid_chunks, dedup_source=True)
            total_chunks += len(valid_chunks)
            total_files += 1

        except Exception as e:
            errors += 1

    elapsed = time.time() - start

    logger.info("\n" + "=" * 60)
    logger.info("📊 RÉSUMÉ")
    logger.info("=" * 60)
    logger.info(f"   Fichiers traités : {total_files}")
    logger.info(f"   Fichiers ignorés : {skipped}")
    logger.info(f"   Erreurs          : {errors}")
    logger.info(f"   Chunks indexés   : {total_chunks}")
    logger.info(f"   Durée            : {elapsed:.1f}s")
    logger.info("=" * 60)
    logger.info("✅ Terminé !")


if __name__ == "__main__":
    main()
